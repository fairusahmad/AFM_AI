"""读取桌面上的 operation.pptx 并打印所有幻灯片内容"""
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation(r'C:\Users\gyixu\Desktop\operation.pptx')

print(f"共 {len(prs.slides)} 页幻灯片\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"{'='*60}")
    print(f"  第 {i} 页")
    print(f"{'='*60}")
    for shape in slide.shapes:
        # 文本框
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f"  [文本] {text}")
        # 表格
        if shape.has_table:
            table = shape.table
            print(f"  [表格] {len(table.rows)} 行 x {len(table.columns)} 列")
            for row_idx, row in enumerate(table.rows):
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                print(f"    Row {row_idx}: {row_text}")
        # 形状名称
        print(f"    (shape: {shape.shape_type}, name: '{shape.name}')")
    print()
