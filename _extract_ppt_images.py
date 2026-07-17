"""提取PPT中所有图片并保存到 _ppt_images/ 目录"""
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import os

prs = Presentation(r'C:\Users\gyixu\Desktop\operation.pptx')
out_dir = Path(r'C:\Users\gyixu\Desktop\AFM-Hysteresis-Simulation-main\_ppt_images')
out_dir.mkdir(exist_ok=True)

image_count = 0
for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            image = shape.image
            ext = image.content_type.split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            fname = f"slide{i}_{shape.name}.{ext}"
            fpath = out_dir / fname
            with open(fpath, 'wb') as f:
                f.write(image.blob)
            image_count += 1
            print(f"Saved: {fname} ({len(image.blob)} bytes)")

print(f"\n共提取 {image_count} 张图片到 {out_dir}")
