from pptx import Presentation

prs = Presentation('operation.pptx')
for i, slide in enumerate(prs.slides):
    print(f'=== 幻灯片 {i+1} ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text)
    print()
